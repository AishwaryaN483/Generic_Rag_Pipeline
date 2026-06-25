# from pptx import Presentation

# def parse_pptx(file_path):
#     prs = Presentation(file_path)
#     text_parts = []
#     table_texts = []

#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text_parts.append(shape.text)
#             if shape.shape_type == 19:  # Table
#                 table = shape.table
#                 rows = []
#                 for r in table.rows:
#                     rows.append(", ".join(cell.text.strip() for cell in r.cells))
#                 table_texts.append("\n".join(rows))

#     combined_text = "\n".join(text_parts)
#     return {"text": combined_text, "tables": table_texts}

from pptx import Presentation
import re


def clean_text(text):
    """
    Clean extracted text:
    - Remove extra spaces
    - Remove weird line breaks
    - Normalize whitespace
    """
    text = text.replace("\n", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def remove_duplicates_preserve_order(text_list):
    seen = set()
    result = []
    for item in text_list:
        if item not in seen and item != "":
            seen.add(item)
            result.append(item)
    return result


def parse_pptx(file_path):
    prs = Presentation(file_path)

    slides_data = []
    full_text_parts = []
    global_seen_text = set()   # 🔥 prevents duplicates across slides

    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_tables = []

        for shape in slide.shapes:

            # -------- TEXT EXTRACTION --------
            if hasattr(shape, "text"):
                raw = shape.text.strip()

                if raw:
                    cleaned = clean_text(raw)

                    # Remove "Slide X -" noise
                    if cleaned.lower().startswith("slide"):
                        parts = cleaned.split("-", 1)
                        if len(parts) > 1:
                            cleaned = parts[1].strip()

                    if cleaned and cleaned not in global_seen_text:
                        slide_text.append(cleaned)
                        global_seen_text.add(cleaned)

            # -------- TABLE EXTRACTION --------
            if shape.has_table:
                table = shape.table
                rows = []

                for r in table.rows:
                    row_text = " | ".join(
                        clean_text(cell.text) for cell in r.cells if cell.text.strip()
                    )
                    if row_text:
                        rows.append(row_text)

                if rows:
                    table_text = "\n".join(rows)

                    if table_text not in global_seen_text:
                        slide_tables.append(table_text)
                        global_seen_text.add(table_text)

        # Remove duplicates inside slide (extra safety)
        slide_text = remove_duplicates_preserve_order(slide_text)

        slide_content = {
            "slide_number": i + 1,
            "text": "\n".join(slide_text),
            "tables": slide_tables
        }

        slides_data.append(slide_content)

        # Build combined text
        if slide_content["text"]:
            full_text_parts.append(f"[Slide {i+1}]\n{slide_content['text']}")

        if slide_tables:
            full_text_parts.append("\n".join(slide_tables))

    combined_text = "\n\n".join(full_text_parts).strip()

    return {                        # metadata+ + content for LLM
        "doc_id": file_path,
        "text": combined_text,
        "slides": slides_data,
        "metadata": {
            "total_slides": len(prs.slides),
            "source": "pptx"
        }
    }


# -------- RUN SCRIPT in Console--------
if __name__ == "__main__":
    file_path = "D://topicTrace//Trail//FSC//Samples//ppts//Sample2.pptx"  # change this

    result = parse_pptx(file_path)

    print("Slides:", result["metadata"]["total_slides"])
    print("\nPreview:\n")
    print(result["text"][:500])

    # Save output for inspection
    with open("parsed_output.txt", "w", encoding="utf-8") as f:
        f.write(result["text"])

    print("\nSaved to parsed_output.txt")